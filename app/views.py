from django.shortcuts import render
from .forms import PowerPointUploadForm


import os
from pptx import Presentation
import boto3
from django.conf import settings

# AWS Keys

aws_access_key_id = settings.AWS_ACCESS_KEY_ID
aws_secret_access_key = settings.AWS_SECRET_ACCESS_KEY 
region_name = 'ap-southeast-1'




def ppt_file_upload(request):
    if request.method == 'POST':
        form = PowerPointUploadForm(request.POST, request.FILES)
        if form.is_valid():
            powerpoint_file = request.FILES['powerpoint_file']

            # Extract notes and generate audio files
            audio_files = extract_notes_and_generate_audio(powerpoint_file)

            # Pass audio files to the template
            return render(request, 'ppt_file_upload.html', {'audio_files': audio_files})
    else:
        form = PowerPointUploadForm()

    return render(request, 'ppt_file_upload.html', {'form': form})


def extract_notes_and_generate_audio(powerpoint_file):
    audio_files = []
    
    # Save the uploaded PowerPoint file temporarily
    with open('temp.pptx', 'wb') as f:
        for chunk in powerpoint_file.chunks():
            f.write(chunk)
    
    # Extract notes from the PowerPoint file
    presentation = Presentation('temp.pptx')
    for i, slide in enumerate(presentation.slides, start=1):
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text
        if notes_text:
            # Generate audio file using Amazon Polly
            audio_file_name = f"slide_{i}_audio.mp3"  # Adjust the file extension as needed
            audio_file_path = generate_audio_file(notes_text, audio_file_name)
            audio_files.append({
                'slide_number': i,
                'audio_filename': audio_file_name,
                'audio_url': audio_file_path.replace(settings.MEDIA_ROOT, settings.MEDIA_URL)  # Adjust the URL path as needed
            })
    
    # Remove the temporary file
    os.remove('temp.pptx')
    
    return audio_files




def generate_audio_file(text, filename):
    client = boto3.client(
        'polly',
        region_name=region_name,  # Replace with your region
        aws_access_key_id=aws_access_key_id,  # Replace with your AWS access key ID
        aws_secret_access_key=aws_secret_access_key  # Replace with your AWS secret access key
    )

    response = client.synthesize_speech(
        Text=text,
        OutputFormat='mp3',
        VoiceId='Raveena'  # Replace with the desired voice ID
    )

    # Ensure the output directory exists
    output_directory = os.path.join(settings.MEDIA_ROOT, 'output_audio')
    os.makedirs(output_directory, exist_ok=True)

    # Save the audio file inside the output directory
    audio_file_path = os.path.join(output_directory, filename)
    with open(audio_file_path, 'wb') as f:
        f.write(response['AudioStream'].read())

    return audio_file_path
        



